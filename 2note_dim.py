import tkinter as tk
from tkinter import ttk, simpledialog, filedialog, colorchooser, messagebox
import tkinter.font as tkfont
from PIL import Image, ImageTk, ImageDraw, ImageFont, ImageGrab, ImageFilter
import platform
import os
import sys
import math
import copy
import traceback
import io
import pickle

# --- FIX LỖI DPI SCALING TRÊN WINDOWS ---
if platform.system() == 'Windows':
    import ctypes
    try:
        ctypes.windll.shcore.SetProcessDpiAwareness(2)
    except Exception:
        try:
            ctypes.windll.user32.SetProcessDPIAware()
        except Exception:
            pass

try:
    import keyboard
    HAS_KEYBOARD = True
except ImportError:
    HAS_KEYBOARD = False

try:
    import cv2
    HAS_CV2 = True
except ImportError:
    HAS_CV2 = False

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

import winreg as reg

# Tương thích PIL mới
RESAMPLE_NEAREST = Image.Resampling.NEAREST if hasattr(Image, 'Resampling') else Image.NEAREST
RESAMPLE_LANCZOS = Image.Resampling.LANCZOS if hasattr(Image, 'Resampling') else Image.LANCZOS
RESAMPLE_BOX = Image.Resampling.BOX if hasattr(Image, 'Resampling') else Image.BOX

# ==========================================
# CÁC HÀM HỖ TRỢ XUẤT FILE THÔNG MINH
# ==========================================
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def add_image_exact(slide, img_stream, left, top, width, height):
    img_stream.seek(0)
    pic = slide.shapes.add_picture(img_stream, int(left), int(top), width=int(width), height=int(height))
    pic.line.color.rgb = RGBColor(30, 30, 30)
    pic.line.width = Inches(0.03) 

def move_slide(prs, old_index, new_index):
    if old_index == new_index: return
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    slide_to_move = slides[old_index]
    xml_slides.remove(slide_to_move)
    xml_slides.insert(new_index, slide_to_move)

def partition_images(imgs, max_size):
    if not imgs: return []
    n = len(imgs)
    num_slides = math.ceil(n / max_size)
    base = n // num_slides
    rem = n % num_slides
    res = []
    idx = 0
    for i in range(num_slides):
        s = base + 1 if i < rem else base
        res.append(imgs[idx:idx+s])
        idx += s
    return res

def draw_adaptive_grid(slide, layout_rows, start_x_base, start_y_base, usable_w, usable_h, GAP):
    if not layout_rows: return
    H_final = (usable_h - (len(layout_rows) - 1) * GAP) / len(layout_rows)
    for row in layout_rows:
        if not row: continue
        sr = sum(img['w'] / img['h'] for img in row)
        H_row_max = (usable_w - (len(row) - 1) * GAP) / sr
        if H_row_max < H_final:
            H_final = H_row_max
            
    Total_H = len(layout_rows) * H_final + (len(layout_rows) - 1) * GAP
    current_y = start_y_base + (usable_h - Total_H) / 2
    
    for row in layout_rows:
        if not row: continue
        row_w = sum(H_final * (img['w'] / img['h']) for img in row) + (len(row) - 1) * GAP
        current_x = start_x_base + (usable_w - row_w) / 2
        for img in row:
            img_w = H_final * (img['w'] / img['h'])
            add_image_exact(slide, img['stream'], current_x, current_y, img_w, H_final)
            current_x += img_w + GAP
        current_y += H_final + GAP

# ==========================================
# LỚP KHUNG CUỘN (SCROLLABLE FRAME)
# ==========================================
class ScrollableFrame(tk.Frame):
    def __init__(self, container, *args, **kwargs):
        super().__init__(container, *args, **kwargs)
        self.canvas = tk.Canvas(self, bg="#252526", highlightthickness=0)
        self.scrollbar = ttk.Scrollbar(self, orient="vertical", command=self.canvas.yview)
        self.scrollable_frame = tk.Frame(self.canvas, bg="#252526")

        self.window_id = self.canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        self.canvas.configure(yscrollcommand=self.scrollbar.set)

        self.canvas.pack(side="left", fill="both", expand=True)
        self.scrollbar.pack(side="right", fill="y")

        self.scrollable_frame.bind("<Configure>", self.on_frame_configure)
        self.canvas.bind("<Configure>", self.on_canvas_configure)
        
        self.bind_mouse_scroll(self.canvas)
        self.bind_mouse_scroll(self.scrollable_frame)

    def on_frame_configure(self, event):
        self.canvas.configure(scrollregion=self.canvas.bbox("all"))

    def on_canvas_configure(self, event):
        self.canvas.itemconfig(self.window_id, width=event.width)

    def bind_mouse_scroll(self, widget):
        widget.bind("<Enter>", self._bind_mouse)
        widget.bind("<Leave>", self._unbind_mouse)

    def _bind_mouse(self, event):
        self.canvas.bind_all("<MouseWheel>", self._on_mousewheel)

    def _unbind_mouse(self, event):
        self.canvas.unbind_all("<MouseWheel>")

    def _on_mousewheel(self, event):
        self.canvas.yview_scroll(int(-1*(event.delta/120)), "units")

# ==========================================
# LỚP HỖ TRỢ CHỤP MÀN HÌNH
# ==========================================
class ScreenSnip:
    def __init__(self, master, mode, callback):
        self.master = master
        self.mode = mode
        self.callback = callback
        self.master.withdraw()
        self.master.update_idletasks()
        self.master.after(400, self.take_snapshot)
        
    def take_snapshot(self):
        self.screen_img = ImageGrab.grab(all_screens=True)
        self.snip_window = tk.Toplevel(self.master)
        self.snip_window.attributes('-fullscreen', True)
        self.snip_window.attributes('-topmost', True)
        self.snip_window.config(cursor="crosshair")
        
        self.canvas = tk.Canvas(self.snip_window, highlightthickness=0)
        self.canvas.pack(fill=tk.BOTH, expand=True)
        
        self.tk_screen = ImageTk.PhotoImage(self.screen_img)
        self.canvas.create_image(0, 0, anchor=tk.NW, image=self.tk_screen)
        
        self.freeform_points = []
        self.rect = None
        self.state = "STARTING"
        self.handle_size = 10
        self.start_x = 0
        self.start_y = 0
        self.drag_start_rect = None
        
        self.canvas.bind("<ButtonPress-1>", self.on_press)
        self.canvas.bind("<B1-Motion>", self.on_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_release)
        self.canvas.bind("<Motion>", self.on_motion)
        self.snip_window.bind("<Return>", lambda e: self.confirm_crop())
        self.snip_window.bind("<Escape>", lambda e: self.cancel())

    def draw_overlay(self):
        if not self.rect or self.mode == "freeform": return
        self.canvas.delete("overlay")
        
        x1 = min(self.rect[0], self.rect[2])
        y1 = min(self.rect[1], self.rect[3])
        x2 = max(self.rect[0], self.rect[2])
        y2 = max(self.rect[1], self.rect[3])
        
        w = self.screen_img.width
        h = self.screen_img.height
        
        self.canvas.create_rectangle(0, 0, w, y1, fill="black", stipple="gray50", outline="", tags="overlay")
        self.canvas.create_rectangle(0, y2, w, h, fill="black", stipple="gray50", outline="", tags="overlay")
        self.canvas.create_rectangle(0, y1, x1, y2, fill="black", stipple="gray50", outline="", tags="overlay")
        self.canvas.create_rectangle(x2, y1, w, y2, fill="black", stipple="gray50", outline="", tags="overlay")
        
        self.canvas.create_rectangle(x1, y1, x2, y2, outline="#007ACC", width=2, tags="overlay")
        
        hs = self.handle_size / 2
        pts = [
            ("TL", x1, y1), ("T", (x1+x2)/2, y1), ("TR", x2, y1),
            ("L", x1, (y1+y2)/2), ("R", x2, (y1+y2)/2),
            ("BL", x1, y2), ("B", (x1+x2)/2, y2), ("BR", x2, y2)
        ]
        
        for name, hx, hy in pts:
            self.canvas.create_rectangle(hx-hs, hy-hs, hx+hs, hy+hs, fill="#007ACC", outline="white", tags=("overlay", f"handle_{name}"))
            
        self.canvas.create_text((x1+x2)/2, y2 + 20, text="Nhấn [ENTER] để chốt | [ESC] để hủy", fill="white", font=("Arial", 12, "bold"), tags="overlay")

    def get_handle(self, x, y):
        if not self.rect: return None
        x1 = min(self.rect[0], self.rect[2])
        y1 = min(self.rect[1], self.rect[3])
        x2 = max(self.rect[0], self.rect[2])
        y2 = max(self.rect[1], self.rect[3])
        
        hs = self.handle_size
        handles = {
            "TL": (x1, y1), "T": ((x1+x2)/2, y1), "TR": (x2, y1),
            "L": (x1, (y1+y2)/2), "R": (x2, (y1+y2)/2),
            "BL": (x1, y2), "B": ((x1+x2)/2, y2), "BR": (x2, y2)
        }
        
        for name, (hx, hy) in handles.items():
            if abs(x - hx) <= hs and abs(y - hy) <= hs:
                return name
                
        if x1 < x < x2 and y1 < y < y2:
            return "CENTER"
        return None

    def on_motion(self, event):
        if self.state in ["STARTING", "DRAWING"] or self.mode == "freeform": return
        h = self.get_handle(event.x, event.y)
        cursors = {
            "TL": "size_nw_se", "BR": "size_nw_se", "TR": "size_ne_sw", "BL": "size_ne_sw",
            "T": "size_ns", "B": "size_ns", "L": "size_we", "R": "size_we",
            "CENTER": "fleur", None: "crosshair"
        }
        self.canvas.config(cursor=cursors.get(h, "crosshair"))

    def on_press(self, event):
        if self.mode == "freeform":
            self.freeform_points = [(event.x, event.y)]
            return
            
        if self.state == "STARTING":
            self.rect = [event.x, event.y, event.x, event.y]
            self.state = "DRAWING"
        else:
            h = self.get_handle(event.x, event.y)
            if h:
                self.state = h if h != "CENTER" else "MOVING"
                self.start_x = event.x
                self.start_y = event.y
                self.drag_start_rect = list(self.rect)
            else:
                self.rect = [event.x, event.y, event.x, event.y]
                self.state = "DRAWING"

    def on_drag(self, event):
        if self.mode == "freeform":
            self.freeform_points.append((event.x, event.y))
            if len(self.freeform_points) > 1:
                self.canvas.create_line(
                    self.freeform_points[-2][0], self.freeform_points[-2][1], 
                    self.freeform_points[-1][0], self.freeform_points[-1][1], 
                    fill="red", width=2, tags="overlay"
                )
            return
            
        if self.state == "DRAWING":
            self.rect[2] = event.x
            self.rect[3] = event.y
        elif self.state == "MOVING":
            dx = event.x - self.start_x
            dy = event.y - self.start_y
            self.rect = [
                self.drag_start_rect[0]+dx, self.drag_start_rect[1]+dy, 
                self.drag_start_rect[2]+dx, self.drag_start_rect[3]+dy
            ]
        elif self.state != "IDLE":
            if "T" in self.state: self.rect[1] = event.y
            if "B" in self.state: self.rect[3] = event.y
            if "L" in self.state: self.rect[0] = event.x
            if "R" in self.state: self.rect[2] = event.x
        self.draw_overlay()

    def on_release(self, event):
        if self.mode == "freeform":
            self.confirm_crop()
        else:
            if self.rect:
                self.rect = [
                    min(self.rect[0], self.rect[2]), min(self.rect[1], self.rect[3]), 
                    max(self.rect[0], self.rect[2]), max(self.rect[1], self.rect[3])
                ]
            self.state = "IDLE"
            self.draw_overlay()

    def confirm_crop(self):
        if self.mode == "freeform":
            if len(self.freeform_points) > 5:
                try:
                    xs = [p[0] for p in self.freeform_points]
                    ys = [p[1] for p in self.freeform_points]
                    bbox = (min(xs), min(ys), max(xs), max(ys))
                    if bbox[2] - bbox[0] > 10 and bbox[3] - bbox[1] > 10:
                        mask = Image.new('L', self.screen_img.size, 0)
                        ImageDraw.Draw(mask).polygon(self.freeform_points, fill=255)
                        res = self.screen_img.copy().convert("RGBA")
                        res.putalpha(mask)
                        self.finish(res.crop(bbox))
                    else:
                        self.cancel()
                except Exception:
                    self.cancel()
            else:
                self.cancel()
        elif self.rect:
            x1, y1, x2, y2 = self.rect
            if x2 - x1 > 10 and y2 - y1 > 10:
                self.finish(self.screen_img.crop((x1, y1, x2, y2)))
            else:
                self.cancel()

    def finish(self, img):
        self.snip_window.destroy()
        self.master.deiconify() 
        self.master.attributes('-topmost', True)
        self.master.after(100, lambda: self.master.attributes('-topmost', False))
        self.callback(img)

    def cancel(self):
        self.snip_window.destroy()
        self.master.deiconify()

# ==========================================
# LỚP CHÍNH ỨNG DỤNG
# ==========================================
class AdvancedNoteDimApp:
    def __init__(self, root):
        self.root = root
        self.root.title("Công cụ Note-Dim V49 (Bản Sạch - Hoàn Hảo Từng Chữ)")
        self.root.geometry("1300x800")
        
        try:
            self.root.state('zoomed')
        except Exception:
            self.root.attributes('-zoomed', True)
            
        self.root.protocol("WM_DELETE_WINDOW", self.exit_app)

        # Khởi tạo các biến hệ thống
        self.raw_images = []
        self.exported_pages = []
        self.left_thumbs = []
        self.right_thumbs = []

        self.original_image = None
        self.tk_image = None
        self.zoom_factor = 1.0
        
        self.undo_stack = []
        self.redo_stack = []
        self.action_counter = 0
        self.current_right_idx = None 
        
        self.current_tool = tk.StringVar(value="dim")
        self.line_color = "#FFD700"
        self.line_width = tk.IntVar(value=10)
        self.dim_unit = tk.StringVar(value="m")
        self.run_on_startup = tk.BooleanVar(value=self.check_startup_registry())
        
        self.font_family = tk.StringVar(value="Arial")
        self.font_size = tk.IntVar(value=15)
        self.text_color = "#000000"
        self.text_bg = "#FFD700"
        self.is_bold = tk.BooleanVar(value=True)
        self.is_italic = tk.BooleanVar(value=False)
        self.is_underline = tk.BooleanVar(value=False)
        self.text_case = tk.StringVar(value="Normal")
        
        self.start_x = None
        self.start_y = None
        self.temp_shapes = []
        self.angle_points = []
        self.draw_points = []
        
        self.selected_action = None
        self.drag_mode = None
        self.drag_start_x = 0
        self.drag_start_y = 0
        self.drag_start_rect = [0, 0, 0, 0]
        self.last_drag_x = 0
        self.last_drag_y = 0
        self.total_dx = 0
        self.total_dy = 0

        self.setup_ui()
        self.bind_shortcuts()
        self.register_global_hotkeys()
        
        self.current_tool.trace_add("write", self.on_tool_change)
        
        if "--background" in sys.argv:
            self.root.withdraw()

    def show_virtual_keyboard(self):
        if platform.system() == 'Windows':
            try:
                os.startfile(r"C:\Program Files\Common Files\microsoft shared\ink\TabTip.exe")
            except Exception:
                pass

    def exit_app(self): 
        try:
            self.root.destroy()
        except Exception:
            pass
        os._exit(0)

    # --- HÀM TẠO DỰ ÁN ---
    def image_to_bytes(self, img, fmt="JPEG"):
        b = io.BytesIO()
        img.save(b, format=fmt, quality=90)
        return b.getvalue()

    def bytes_to_image(self, b_data):
        return Image.open(io.BytesIO(b_data))

    def save_project(self):
        if not self.raw_images and not self.exported_pages:
            messagebox.showwarning("Trống", "Không có dữ liệu để lưu dự án!")
            return
            
        fp = filedialog.asksaveasfilename(defaultextension=".ndp", filetypes=[("Dự Án Note-Dim", "*.ndp")], initialfile="MyProject.ndp")
        if not fp: return
        
        try:
            dt = {'raws': [], 'exports': []}
            for img_dict in self.raw_images:
                dt['raws'].append(self.image_to_bytes(img_dict['image'].convert('RGB')))
                
            for pg in self.exported_pages:
                clean_stack = []
                for act in pg['undo_stack']:
                    new_act = {k: v for k, v in act.items() if k != 'tk_cache'}
                    if new_act.get('img_data'):
                        new_act['img_data'] = self.image_to_bytes(new_act['img_data'], "PNG")
                    clean_stack.append(new_act)
                
                dt['exports'].append({
                    'raw_img_bytes': self.image_to_bytes(pg['raw_image'].convert('RGB')),
                    'stack': clean_stack
                })
                
            with open(fp, 'wb') as f:
                pickle.dump(dt, f)
            messagebox.showinfo("Thành công", "Đã lưu Dự án an toàn!")
        except Exception as e:
            messagebox.showerror("Lỗi", f"Không thể lưu dự án:\n{e}")

    def open_project(self):
        fp = filedialog.askopenfilename(filetypes=[("Dự Án Note-Dim", "*.ndp")])
        if not fp: return
        
        try:
            with open(fp, 'rb') as f:
                dt = pickle.load(f)
                
            self.raw_images = [{'image': self.bytes_to_image(b).convert('RGB')} for b in dt['raws']]
            self.exported_pages = []
            
            for pg in dt['exports']:
                restored_stack = []
                for act in pg['stack']:
                    if act.get('img_data') and isinstance(act['img_data'], bytes):
                        act['img_data'] = self.bytes_to_image(act['img_data']).convert('RGBA')
                    restored_stack.append(act)
                    
                ri = self.bytes_to_image(pg['raw_img_bytes']).convert('RGB')
                
                backup_orig = self.original_image
                backup_stack = self.undo_stack
                
                self.original_image = ri
                self.undo_stack = restored_stack
                ei = self.get_export_image()
                
                self.original_image = backup_orig
                self.undo_stack = backup_stack
                
                self.exported_pages.append({'raw_image': ri, 'undo_stack': restored_stack, 'rendered_image': ei})
            
            self.refresh_left_panel()
            self.refresh_right_panel()
            
            if self.raw_images:
                self.load_from_left(0)
                
            messagebox.showinfo("Thành công", "Đã mở Dự án thành công!")
        except Exception as e:
            messagebox.showerror("Lỗi", f"Không thể mở dự án:\n{e}")

    # --- HÀM KHỞI ĐỘNG VÀ HOTKEY ---
    def check_startup_registry(self):
        if platform.system() != 'Windows': return False
        try:
            key = reg.OpenKey(reg.HKEY_CURRENT_USER, r"Software\Microsoft\Windows\CurrentVersion\Run", 0, reg.KEY_READ)
            reg.QueryValueEx(key, "NoteDimApp")
            reg.CloseKey(key)
            return True
        except WindowsError:
            return False

    def toggle_startup(self):
        if platform.system() != 'Windows': return
        key_path = r"Software\Microsoft\Windows\CurrentVersion\Run"
        app_name = "NoteDimApp"
        try:
            key = reg.OpenKey(reg.HKEY_CURRENT_USER, key_path, 0, reg.KEY_ALL_ACCESS)
            if self.run_on_startup.get():
                cmd = f'"{sys.executable}" "{os.path.abspath(sys.argv[0])}" --background' if sys.argv[0].endswith('.py') else f'"{os.path.abspath(sys.argv[0])}" --background'
                reg.SetValueEx(key, app_name, 0, reg.REG_SZ, cmd)
            else:
                reg.DeleteValue(key, app_name)
            reg.CloseKey(key)
        except Exception:
            pass

    def register_global_hotkeys(self):
        if HAS_KEYBOARD:
            keyboard.add_hotkey('ctrl+alt+q', lambda: self.root.after(0, lambda: self.start_snip("rectangle")))

    # --- HÀM CHỤP ẢNH ---
    def start_snip(self, mode):
        if mode == "fullscreen":
            self.root.withdraw()
            self.root.after(300, lambda: self.load_captured_image(ImageGrab.grab(all_screens=True)))
        else:
            ScreenSnip(self.root, mode, self.load_captured_image)

    def load_captured_image(self, img):
        if self.root.state() != 'normal':
            self.root.deiconify()
            
        if img.mode == 'RGBA':
            bg = Image.new("RGB", img.size, (255, 255, 255))
            bg.paste(img, mask=img.split()[3])
            img = bg
            
        self.raw_images.append({'image': img})
        self.refresh_left_panel()
        self.load_from_left(len(self.raw_images) - 1)

    # --- HÀM CAMERA ---
    def open_camera(self):
        if not HAS_CV2:
            messagebox.showwarning("Lỗi", "Vui lòng mở CMD chạy:\npip install opencv-python")
            return
            
        self.cam_window = tk.Toplevel(self.root)
        self.cam_window.title("Chụp Camera")
        self.cam_window.geometry("800x600")
        self.cam_window.transient(self.root)
        self.cam_window.grab_set()
        
        self.cam_label = tk.Label(self.cam_window, bg="black")
        self.cam_label.pack(fill=tk.BOTH, expand=True)
        
        btn_frame = tk.Frame(self.cam_window, bg="#2D2D30", height=60)
        btn_frame.pack(fill=tk.X, side=tk.BOTTOM)
        
        tk.Button(btn_frame, text="🔄 Đổi Camera", command=self.switch_camera, bg="#17A2B8", fg="white", font=("Arial", 12, "bold")).pack(side=tk.LEFT, padx=10, pady=10)
        tk.Button(btn_frame, text="📸 CHỤP", command=self.take_cam_snapshot, bg="#28A745", fg="white", font=("Arial", 12, "bold")).pack(side=tk.LEFT, expand=True, padx=10, pady=10)
        tk.Button(btn_frame, text="❌ Đóng", command=self.close_camera, bg="#DC3545", fg="white", font=("Arial", 12, "bold")).pack(side=tk.RIGHT, padx=10, pady=10)
        
        self.cam_idx = 1
        self.cap = cv2.VideoCapture(self.cam_idx)
        if not self.cap.isOpened():
            self.cam_idx = 0
            self.cap = cv2.VideoCapture(self.cam_idx)
            
        self.cam_window.protocol("WM_DELETE_WINDOW", self.close_camera)
        self.update_camera()

    def update_camera(self):
        if hasattr(self, 'cap') and self.cap.isOpened():
            ret, frame = self.cap.read()
            if ret:
                img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                ww = self.cam_window.winfo_width()
                wh = self.cam_window.winfo_height() - 60
                if ww > 10 and wh > 10:
                    img.thumbnail((ww, wh))
                tk_img = ImageTk.PhotoImage(image=img)
                self.cam_label.imgtk = tk_img
                self.cam_label.configure(image=tk_img)
            self.cam_job = self.cam_window.after(30, self.update_camera)

    def switch_camera(self):
        if hasattr(self, 'cap'):
            self.cap.release()
        self.cam_idx = 1 - self.cam_idx
        self.cap = cv2.VideoCapture(self.cam_idx)
        if not self.cap.isOpened():
            self.cam_idx = 1 - self.cam_idx
            self.cap = cv2.VideoCapture(self.cam_idx)

    def take_cam_snapshot(self):
        if hasattr(self, 'cap') and self.cap.isOpened():
            ret, frame = self.cap.read()
            if ret:
                img = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                self.close_camera()
                self.load_captured_image(img)
        else:
            self.close_camera()

    def close_camera(self):
        if hasattr(self, 'cam_job'):
            self.cam_window.after_cancel(self.cam_job)
        if hasattr(self, 'cap'):
            self.cap.release()
        self.cam_window.destroy()

    # --- HÀM ZOOM ---
    def zoom_in(self):
        if not self.original_image: return
        self.zoom_factor *= 1.1
        self.apply_zoom()

    def zoom_out(self):
        if not self.original_image: return
        self.zoom_factor *= 0.9
        self.apply_zoom()

    def zoom_reset(self):
        if not self.original_image: return
        self.zoom_factor = 1.0
        self.apply_zoom()

    def apply_zoom(self):
        self.update_image_display()
        self.redraw_all()

    # --- GIAO DIỆN CHÍNH ---
    def setup_ui(self):
        # Toolbar 1
        t1 = tk.Frame(self.root, bg="#2D2D30", height=40)
        t1.pack(side=tk.TOP, fill=tk.X)
        
        prj_btn = tk.Menubutton(t1, text="📁 Dự Án ▼", bg="#17A2B8", fg="white", font=("Arial", 9, "bold"), relief="raised")
        prj_btn.pack(side=tk.LEFT, padx=5, pady=5)
        pm = tk.Menu(prj_btn, tearoff=0)
        pm.add_command(label="📂 Mở Dự Án (.ndp)", command=self.open_project)
        pm.add_command(label="💾 Lưu Dự Án (.ndp)", command=self.save_project)
        prj_btn["menu"] = pm

        tk.Button(t1, text="➕ Thêm Ảnh", command=self.add_multiple_images, bg="#007ACC", fg="white", font=("Arial", 9, "bold")).pack(side=tk.LEFT, padx=5)
        
        bs = tk.Menubutton(t1, text="📸 Chụp", bg="#E0A800", fg="black", font=("Arial", 9, "bold"), relief="raised")
        bs.pack(side=tk.LEFT, padx=5)
        sm = tk.Menu(bs, tearoff=0)
        sm.add_command(label="Rectangle", command=lambda: self.start_snip("rectangle"))
        sm.add_command(label="Freeform", command=lambda: self.start_snip("freeform"))
        bs["menu"] = sm
        
        tk.Button(t1, text="📷 Cam", command=self.open_camera, bg="#6F42C1", fg="white", font=("Arial", 9, "bold")).pack(side=tk.LEFT, padx=5)
        tk.Button(t1, text="📥 Lưu DS", command=self.add_to_right_panel, bg="#28A745", fg="white", font=("Arial", 9, "bold")).pack(side=tk.LEFT, padx=5)
        
        tk.Label(t1, text=" | ", bg="#2D2D30", fg="white").pack(side=tk.LEFT)
        tk.Button(t1, text="-", command=self.zoom_out, bg="#555", fg="white").pack(side=tk.LEFT, padx=2)
        tk.Button(t1, text="1:1", command=self.zoom_reset, bg="#555", fg="white").pack(side=tk.LEFT, padx=2)
        tk.Button(t1, text="+", command=self.zoom_in, bg="#555", fg="white").pack(side=tk.LEFT, padx=2)
        tk.Label(t1, text=" | ", bg="#2D2D30", fg="white").pack(side=tk.LEFT)
        
        tk.Button(t1, text="↶ Undo", command=self.undo, bg="#555", fg="white").pack(side=tk.LEFT, padx=10)
        tk.Button(t1, text="↷ Redo", command=self.redo, bg="#555", fg="white").pack(side=tk.LEFT, padx=2)
        
        tk.Button(t1, text="❌ Thoát", command=self.exit_app, bg="#DC3545", fg="white", font=("Arial", 9, "bold")).pack(side=tk.RIGHT, padx=5, pady=5)

        # Toolbar 2
        t2 = tk.Frame(self.root, bg="#3E3E42", height=40)
        t2.pack(side=tk.TOP, fill=tk.X)
        
        tk.Radiobutton(t2, text="✋ Kéo", variable=self.current_tool, value="pan", bg="#3E3E42", fg="#00FF00", selectcolor="#555", font=("Arial", 9, "bold")).pack(side=tk.LEFT)
        tk.Radiobutton(t2, text="🖱️ Sửa", variable=self.current_tool, value="select", bg="#3E3E42", fg="white", selectcolor="#555").pack(side=tk.LEFT)
        
        self.draw_btn = tk.Menubutton(t2, text="🎨 Hình Vẽ ▼", bg="#3E3E42", fg="#FFD700", relief="raised", font=("Arial", 9, "bold"))
        self.draw_btn.pack(side=tk.LEFT, padx=2)
        dm = tk.Menu(self.draw_btn, tearoff=0)
        
        def set_draw_tool(val, name):
            self.current_tool.set(val)
            self.draw_btn.config(text=f"{name} ▼")
            
        dm.add_command(label="✍️ Tự do", command=lambda: set_draw_tool("draw", "✍️ Tự do"))
        dm.add_command(label="➖ Thẳng", command=lambda: set_draw_tool("line", "➖ Thẳng"))
        dm.add_command(label="⬜ Vuông", command=lambda: set_draw_tool("rect", "⬜ Vuông"))
        dm.add_command(label="⭕ Tròn", command=lambda: set_draw_tool("oval", "⭕ Tròn"))
        dm.add_separator()
        dm.add_command(label="🌫️ Khảm", command=lambda: set_draw_tool("blur", "🌫️ Khảm"))
        dm.add_command(label="🖼️ Chèn Ảnh", command=self.action_insert_image)
        self.draw_btn["menu"] = dm
        
        tk.Radiobutton(t2, text="📏 Dim", variable=self.current_tool, value="dim", bg="#3E3E42", fg="white", selectcolor="#555").pack(side=tk.LEFT)
        tk.Radiobutton(t2, text="📝 Tag", variable=self.current_tool, value="tag", bg="#3E3E42", fg="white", selectcolor="#555").pack(side=tk.LEFT)
        tk.Radiobutton(t2, text="📐 Góc", variable=self.current_tool, value="angle", bg="#3E3E42", fg="white", selectcolor="#555").pack(side=tk.LEFT)
        tk.Radiobutton(t2, text="↗ Mũi tên", variable=self.current_tool, value="arrow_note", bg="#3E3E42", fg="white", selectcolor="#555").pack(side=tk.LEFT)
        
        tk.Label(t2, text=" | Màu:", bg="#3E3E42", fg="white").pack(side=tk.LEFT, padx=(10,2))
        self.btn_color = tk.Button(t2, text="   ", command=self.choose_line_color, bg=self.line_color, width=3)
        self.btn_color.pack(side=tk.LEFT, padx=2)
        
        tk.Label(t2, text="Nét:", bg="#3E3E42", fg="white").pack(side=tk.LEFT, padx=(5,2))
        tk.Spinbox(t2, from_=1, to=50, textvariable=self.line_width, width=4).pack(side=tk.LEFT, padx=2)
        
        tk.Label(t2, text="Đơn Vị:", bg="#3E3E42", fg="white").pack(side=tk.LEFT, padx=(10,2))
        ttk.Combobox(t2, textvariable=self.dim_unit, values=["m", "mm", "cm", "inch", ""], width=5).pack(side=tk.LEFT, padx=2)

        # Toolbar 3
        t3 = tk.Frame(self.root, bg="#2D2D30", height=40)
        t3.pack(side=tk.TOP, fill=tk.X)
        
        tk.Label(t3, text="Font:", bg="#2D2D30", fg="white", font=("Arial", 9, "bold")).pack(side=tk.LEFT, padx=5)
        ttk.Combobox(t3, textvariable=self.font_family, values=["Arial", "Times New Roman", "Courier New", "Tahoma"], width=15, state="readonly").pack(side=tk.LEFT, padx=2)
        
        tk.Spinbox(t3, from_=8, to=200, textvariable=self.font_size, width=4).pack(side=tk.LEFT, padx=2)
        tk.Checkbutton(t3, text="B", variable=self.is_bold, font=("Arial", 9, "bold"), indicatoron=False, width=2).pack(side=tk.LEFT, padx=1)
        tk.Checkbutton(t3, text="I", variable=self.is_italic, font=("Arial", 9, "italic"), indicatoron=False, width=2).pack(side=tk.LEFT, padx=1)
        tk.Checkbutton(t3, text="U", variable=self.is_underline, font=("Arial", 9, "underline"), indicatoron=False, width=2).pack(side=tk.LEFT, padx=1)
        
        tk.Label(t3, text="Chữ:", bg="#2D2D30", fg="white").pack(side=tk.LEFT, padx=(10,2))
        self.btn_text_color = tk.Button(t3, text="  ", command=self.choose_text_color, bg=self.text_color, fg="white", width=2)
        self.btn_text_color.pack(side=tk.LEFT)
        
        tk.Label(t3, text="Nền:", bg="#2D2D30", fg="white").pack(side=tk.LEFT, padx=(10,2))
        self.btn_text_bg = tk.Button(t3, text="  ", command=self.choose_text_bg, bg=self.text_bg, fg="black", width=2)
        self.btn_text_bg.pack(side=tk.LEFT)
        
        tk.Button(t3, text="✨ Đồng bộ (Ctrl+A)", command=self.batch_edit_dialog, bg="#007ACC", fg="white", font=("Arial", 9, "bold")).pack(side=tk.RIGHT, padx=15)

        # Main Body
        self.main_body = tk.Frame(self.root, bg="#1E1E1E")
        self.main_body.pack(fill=tk.BOTH, expand=True)

        self.left_pane = tk.Frame(self.main_body, width=280, bg="#252526")
        self.left_pane.pack(side=tk.LEFT, fill=tk.Y)
        self.left_pane.pack_propagate(False)
        tk.Label(self.left_pane, text="ẢNH GỐC", bg="#333", fg="white", font=("Arial", 10, "bold"), pady=5).pack(fill=tk.X)
        self.left_scroll = ScrollableFrame(self.left_pane)
        self.left_scroll.pack(fill=tk.BOTH, expand=True)
        self.refresh_left_panel()

        self.right_pane = tk.Frame(self.main_body, width=280, bg="#252526")
        self.right_pane.pack(side=tk.RIGHT, fill=tk.Y)
        self.right_pane.pack_propagate(False)
        tk.Label(self.right_pane, text="DANH SÁCH", bg="#333", fg="#28A745", font=("Arial", 10, "bold"), pady=5).pack(fill=tk.X)
        self.right_scroll = ScrollableFrame(self.right_pane)
        self.right_scroll.pack(side=tk.TOP, fill=tk.BOTH, expand=True)
        
        self.right_bottom = tk.Frame(self.right_pane, bg="#2D2D30")
        self.right_bottom.pack(side=tk.BOTTOM, fill=tk.X)
        tk.Button(self.right_bottom, text="🔨 KHẮC XUẤT", command=self.open_export_dialog, bg="#6F42C1", fg="white", font=("Arial", 11, "bold")).pack(fill=tk.X, padx=5, pady=10)
        tk.Button(self.right_bottom, text="🗑️ Xóa DS", command=self.clear_export_list, bg="#333", fg="white", font=("Arial", 9)).pack(fill=tk.X, padx=5, pady=5)

        self.center_pane = tk.Frame(self.main_body, bg="#1E1E1E")
        self.center_pane.pack(fill=tk.BOTH, expand=True)
        self.canvas = tk.Canvas(self.center_pane, bg="#1E1E1E", cursor="crosshair")
        self.canvas.pack(fill=tk.BOTH, expand=True)
        
        self.canvas.bind("<ButtonPress-1>", self.on_left_press)
        self.canvas.bind("<B1-Motion>", self.on_left_drag)
        self.canvas.bind("<ButtonRelease-1>", self.on_left_release)
        self.canvas.bind("<Double-Button-1>", self.on_double_click) 
        self.canvas.bind("<ButtonPress-3>", self.on_right_press)
        self.canvas.bind("<B3-Motion>", self.on_right_drag)

    # --- KHẮC XUẤT ---
    def clear_export_list(self):
        if messagebox.askyesno("Xác nhận", "Xóa toàn bộ danh sách?"):
            self.exported_pages = []
            self.current_right_idx = None
            self.refresh_right_panel()

    def open_export_dialog(self):
        if not self.exported_pages:
            messagebox.showwarning("Trống", "Chưa có ảnh nào để xuất.")
            return

        self.show_virtual_keyboard()
        d = tk.Toplevel(self.root)
        d.title("Tùy chọn Khắc Xuất Thông Minh")
        d.geometry("900x650")
        d.transient(self.root)
        d.grab_set()

        tk.Label(d, text="⚙️ CÀI ĐẶT DÀN TRANG & XUẤT", font=("Arial", 13, "bold"), fg="#007ACC").pack(pady=(10, 5))
        
        nb = ttk.Notebook(d)
        nb.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
        
        # TAB 1: CÀI ĐẶT CHUNG
        f1 = tk.Frame(nb, bg="#f0f0f0")
        nb.add(f1, text="Cài Đặt Chung & Phạm Vi")
        
        tk.Label(f1, text="Tên file:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=0, column=0, padx=10, pady=15, sticky="w")
        fn_var = tk.StringVar(value="Report")
        tk.Entry(f1, textvariable=fn_var, width=40).grid(row=0, column=1, sticky="w", pady=15)
        
        pm_var = tk.IntVar(value=2)
        tk.Label(f1, text="Chế độ:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=1, column=0, padx=10, pady=5, sticky="w")
        tk.Radiobutton(f1, text="Dùng Template Mẫu", variable=pm_var, value=1, bg="#f0f0f0").grid(row=1, column=1, sticky="w")
        tk.Radiobutton(f1, text="Tạo Mới 100%", variable=pm_var, value=2, bg="#f0f0f0").grid(row=1, column=1, sticky="e")
        
        tk.Label(f1, text="Mẫu PPTX (.pptx):", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=2, column=0, padx=10, pady=5, sticky="w")
        tp_var = tk.StringVar()
        tk.Entry(f1, textvariable=tp_var, width=40).grid(row=2, column=1, sticky="w")
        tk.Button(f1, text="Duyệt...", command=lambda: tp_var.set(filedialog.askopenfilename(filetypes=[("PPTX", "*.pptx")]) or tp_var.get())).grid(row=2, column=2, padx=5)
        
        tk.Label(f1, text="Chèn sau slide số:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=3, column=0, padx=10, pady=5, sticky="w")
        ip_var = tk.IntVar(value=0)
        tk.Spinbox(f1, from_=0, to=999, textvariable=ip_var, width=5).grid(row=3, column=1, sticky="w")
        tk.Label(f1, text="(0 = chèn vào cuối cùng)", bg="#f0f0f0", fg="gray").grid(row=3, column=1, sticky="e")
        
        tk.Label(f1, text="Phạm vi trang:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=4, column=0, padx=10, pady=20, sticky="w")
        md_var = tk.IntVar(value=1)
        fr_opts = tk.Frame(f1, bg="#f0f0f0")
        fr_opts.grid(row=4, column=1, sticky="w", pady=20)
        
        def omc():
            v = md_var.get()
            sp.config(state="normal" if v==2 else "disabled")
            ss.config(state="normal" if v==3 else "disabled")
            se.config(state="normal" if v==3 else "disabled")
            
        tk.Radiobutton(fr_opts, text="Tất cả các trang", variable=md_var, value=1, command=omc, bg="#f0f0f0").grid(row=0, column=0, columnspan=3, sticky="w", pady=2)
        tk.Radiobutton(fr_opts, text="Trang cụ thể:", variable=md_var, value=2, command=omc, bg="#f0f0f0").grid(row=1, column=0, sticky="w", pady=2)
        
        p_var = tk.IntVar(value=1)
        sp = tk.Spinbox(fr_opts, from_=1, to=len(self.exported_pages), textvariable=p_var, width=5, state="disabled")
        sp.grid(row=1, column=1, sticky="w")
        
        tk.Radiobutton(fr_opts, text="Từ trang:", variable=md_var, value=3, command=omc, bg="#f0f0f0").grid(row=2, column=0, sticky="w", pady=2)
        
        frg = tk.Frame(fr_opts, bg="#f0f0f0")
        frg.grid(row=2, column=1, sticky="w")
        
        s_var = tk.IntVar(value=1)
        ss = tk.Spinbox(frg, from_=1, to=len(self.exported_pages), textvariable=s_var, width=4, state="disabled")
        ss.pack(side=tk.LEFT)
        
        tk.Label(frg, text=" đến ", bg="#f0f0f0").pack(side=tk.LEFT)
        
        e_var = tk.IntVar(value=len(self.exported_pages))
        se = tk.Spinbox(frg, from_=1, to=len(self.exported_pages), textvariable=e_var, width=4, state="disabled")
        se.pack(side=tk.LEFT)

        # TAB 2: TRANG GIỮA
        f2 = tk.Frame(nb, bg="#f0f0f0")
        nb.add(f2, text="Trang Giữa (Nội Dung)")
        
        tk.Label(f2, text="Dàn trang:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=0, column=0, padx=10, pady=15, sticky="w")
        lv = tk.StringVar(value="L1")
        tk.Radiobutton(f2, text="L1 (1-2 ảnh)", variable=lv, value="L1", bg="#f0f0f0").grid(row=0, column=1, sticky="w")
        tk.Radiobutton(f2, text="L2 (Lưới thông minh)", variable=lv, value="L2", bg="#f0f0f0").grid(row=0, column=1, sticky="e")
        
        tk.Label(f2, text="Căn lề L1:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=1, column=0, padx=10, pady=5, sticky="w")
        av = tk.StringVar(value="2")
        fa = tk.Frame(f2, bg="#f0f0f0")
        fa.grid(row=1, column=1, sticky="w")
        tk.Radiobutton(fa, text="Trái", variable=av, value="1", bg="#f0f0f0").pack(side=tk.LEFT)
        tk.Radiobutton(fa, text="Giữa", variable=av, value="2", bg="#f0f0f0").pack(side=tk.LEFT, padx=10)
        tk.Radiobutton(fa, text="Phải", variable=av, value="3", bg="#f0f0f0").pack(side=tk.LEFT)
        
        tk.Label(f2, text="Tiêu đề trang (Phụ):", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=2, column=0, padx=10, pady=15, sticky="w")
        t2_var = tk.StringVar(value="HÌNH ẢNH THỰC TẾ")
        tk.Entry(f2, textvariable=t2_var, width=40).grid(row=2, column=1, sticky="w")
        
        tk.Label(f2, text="Ảnh Nền (Trang giữa):", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=3, column=0, padx=10, pady=5, sticky="w")
        mbg_var = tk.StringVar()
        tk.Entry(f2, textvariable=mbg_var, width=40).grid(row=3, column=1, sticky="w")
        tk.Button(f2, text="Duyệt...", command=lambda: mbg_var.set(filedialog.askopenfilename(filetypes=[("Image", "*.jpg;*.png")]) or mbg_var.get())).grid(row=3, column=2, padx=5)

        # TAB 3: BÌA & KẾT
        f3 = tk.Frame(nb, bg="#f0f0f0")
        nb.add(f3, text="Trang Bìa & Kết (Tạo Mới)")
        
        tk.Label(f3, text="Tiêu đề Bìa CHÍNH:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=0, column=0, padx=10, pady=15, sticky="w")
        c_tit = tk.StringVar(value="BÁO CÁO CÔNG VIỆC")
        tk.Entry(f3, textvariable=c_tit, width=40).grid(row=0, column=1, sticky="w")
        
        tk.Label(f3, text="Tiêu đề Phụ:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=1, column=0, padx=10, pady=5, sticky="w")
        c_sub = tk.StringVar(value="Cập nhật tiến độ dự án")
        tk.Entry(f3, textvariable=c_sub, width=40).grid(row=1, column=1, sticky="w")
        
        tk.Label(f3, text="Chữ Trang Kết:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=2, column=0, padx=10, pady=15, sticky="w")
        end_tit = tk.StringVar(value="THANK YOU!")
        tk.Entry(f3, textvariable=end_tit, width=40).grid(row=2, column=1, sticky="w")
        
        tk.Label(f3, text="Màu chữ:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=3, column=0, padx=10, pady=5, sticky="w")
        c_col = tk.StringVar(value="#003366")
        bc_c = tk.Button(f3, bg=c_col.get(), width=5, command=lambda: [c_col.set(colorchooser.askcolor()[1] or c_col.get()), bc_c.config(bg=c_col.get())])
        bc_c.grid(row=3, column=1, sticky="w")
        
        tk.Label(f3, text="Vị trí chữ (1-8):", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=4, column=0, padx=10, pady=15, sticky="w")
        c_pos = tk.StringVar(value="7")
        cpp = tk.Frame(f3, bg="#f0f0f0")
        cpp.grid(row=4, column=1, sticky="w")
        pos_names = ["Trái Trên", "Giữa Trên", "Phải Trên", "Trái Giữa", "Phải Giữa", "Trái Dưới", "Giữa Dưới", "Phải Dưới"]
        for i, pt in enumerate(pos_names):
            tk.Radiobutton(cpp, text=str(i+1), variable=c_pos, value=str(i+1), bg="#f0f0f0").grid(row=i//3, column=i%3, padx=5)
            
        tk.Label(f3, text="Ảnh Nền Bìa:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=5, column=0, padx=10, pady=5, sticky="w")
        c_bg = tk.StringVar()
        tk.Entry(f3, textvariable=c_bg, width=40).grid(row=5, column=1, sticky="w")
        tk.Button(f3, text="Duyệt...", command=lambda: c_bg.set(filedialog.askopenfilename(filetypes=[("Image", "*.jpg;*.png")]) or c_bg.get())).grid(row=5, column=2, padx=5)
        
        tk.Label(f3, text="Ảnh Nền Kết:", font=("Arial", 10, "bold"), bg="#f0f0f0").grid(row=6, column=0, padx=10, pady=15, sticky="w")
        e_bg = tk.StringVar()
        tk.Entry(f3, textvariable=e_bg, width=40).grid(row=6, column=1, sticky="w")
        tk.Button(f3, text="Duyệt...", command=lambda: e_bg.set(filedialog.askopenfilename(filetypes=[("Image", "*.jpg;*.png")]) or e_bg.get())).grid(row=6, column=2, padx=5)

        def get_pgs():
            tot = len(self.exported_pages)
            m = md_var.get()
            fn = fn_var.get()
            if m == 1:
                return self.exported_pages, fn
            elif m == 2:
                if 1 <= p_var.get() <= tot:
                    return [self.exported_pages[p_var.get()-1]], fn
                return [], fn
            else:
                if 1 <= s_var.get() <= e_var.get() <= tot:
                    return self.exported_pages[s_var.get()-1:e_var.get()], fn
                return [], fn

        def exp_ppt():
            if not HAS_PPTX:
                messagebox.showwarning("Lỗi", "Vui lòng cài đặt: pip install python-pptx")
                return
            pgs, fn = get_pgs()
            if not pgs:
                messagebox.showwarning("Lỗi", "Phạm vi trang lỗi!")
                return
                
            sp = filedialog.asksaveasfilename(defaultextension=".pptx", filetypes=[("PowerPoint", "*.pptx")], initialfile=f"{fn}.pptx")
            if not sp: return
            
            try:
                isp = pm_var.get() == 1
                if isp:
                    if not tp_var.get() or not os.path.exists(tp_var.get()):
                        return messagebox.showwarning("Lỗi", "Vui lòng chọn file mẫu .pptx hợp lệ!")
                    prs = Presentation(tp_var.get())
                    vt = ip_v.get()
                    if vt <= 0 or vt > len(prs.slides):
                        vt = len(prs.slides)
                else:
                    prs = Presentation()
                    prs.slide_width = Inches(13.333)
                    prs.slide_height = Inches(7.5)
                    vt = 0
                    
                sw = prs.slide_width
                sh = prs.slide_height
                
                try: sly = prs.slide_layouts[6]
                except: sly = prs.slide_layouts[0]
                
                # --- TRANG BÌA ---
                def add_cover(is_end=False):
                    sl = prs.slides.add_slide(sly)
                    bg_path = e_bg.get() if is_end else c_bg.get()
                    
                    if bg_path and os.path.exists(bg_path):
                        try:
                            sl.shapes.add_picture(bg_path, 0, 0, sw, sh)
                        except:
                            pass
                            
                    ps = c_pos.get()
                    tw = Inches(10)
                    th = Inches(2)
                    
                    cx, cy = {
                        "1":(Inches(0.5), Inches(0.5)), 
                        "2":((sw-tw)/2, Inches(0.5)), 
                        "3":(sw-tw-Inches(0.5), Inches(0.5)), 
                        "4":(Inches(0.5), (sh-th)/2), 
                        "5":(sw-tw-Inches(0.5), (sh-th)/2), 
                        "6":(Inches(0.5), sh-th-Inches(0.5)), 
                        "7":((sw-tw)/2, sh-th-Inches(0.5)), 
                        "8":(sw-tw-Inches(0.5), sh-th-Inches(0.5))
                    }.get(ps, ((sw-tw)/2, sh-th-Inches(0.5)))
                    
                    tx = sl.shapes.add_textbox(cx, cy, tw, th)
                    tf = tx.text_frame
                    p = tf.paragraphs[0]
                    p.text = end_tit.get().upper() if is_end else c_tit.get().upper()
                    p.font.size = Pt(44)
                    p.font.bold = True
                    try:
                        p.font.color.rgb = RGBColor(*hex_to_rgb(c_col.get()))
                    except:
                        pass
                    
                    if not is_end and c_sub.get():
                        p2 = tf.add_paragraph()
                        p2.text = c_sub.get()
                        p2.font.size = Pt(28)
                        try:
                            p2.font.color.rgb = RGBColor(*hex_to_rgb(c_col.get()))
                        except:
                            pass
                    return sl

                # Trang bìa mở đầu
                if not isp and (c_tit.get() or c_bg.get()):
                    add_cover(is_end=False)
                    move_slide(prs, len(prs.slides)-1, 0)
                    vt += 1
                
                imd = []
                for p in pgs:
                    im = p['rendered_image'].convert("RGB")
                    stm = io.BytesIO()
                    im.save(stm, format="JPEG", quality=90)
                    imd.append({'stream': stm, 'w': im.width, 'h': im.height, 'is_portrait': im.height >= im.width})
                
                uw = sw - Inches(0.4)
                uh = sh - Inches(1.7) if (not isp and t2_var.get()) else sh - Inches(1.0)
                G = Inches(0.12)
                CL = Inches(0.2)
                CT = Inches(0.9) if (not isp and t2_var.get()) else Inches(0.5)
                
                l_mode = "L1" if isp else lv.get()
                
                if l_mode == "L1":
                    i = 0
                    while i < len(imd):
                        ci = imd[i]
                        sl = prs.slides.add_slide(sly)
                        
                        if not isp and mbg_var.get() and os.path.exists(mbg_var.get()):
                            try:
                                sl.shapes.add_picture(mbg_var.get(), 0, 0, sw, sh)
                            except:
                                pass
                                
                        if not isp and t2_var.get():
                            p = sl.shapes.add_textbox(Inches(0.2), Inches(0.15), sw-Inches(0.4), Inches(0.6)).text_frame.paragraphs[0]
                            p.text = t2_var.get().upper()
                            p.font.size = Pt(22)
                            p.font.bold = True
                            p.font.underline = True
                            p.font.color.rgb = RGBColor(0, 51, 102)
                            
                        move_slide(prs, len(prs.slides)-1, vt)
                        
                        if ci['is_portrait'] and i+1 < len(imd) and imd[i+1]['is_portrait']:
                            ni = imd[i+1]
                            r1 = ci['w']/ci['h']
                            r2 = ni['w']/ni['h']
                            fh = uh if uh*r1+uh*r2 <= uw-G else (uw-G)/(r1+r2)
                            fw1 = fh*r1
                            fw2 = fh*r2
                            bw = fw1 + G + fw2
                            
                            align = av.get() if not isp else '2'

                            if align == '1': sx = CL
                            elif align == '3': sx = sw - Inches(0.2) - bw
                            else: sx = CL + (uw - bw)/2
                                
                            sy = CT + (uh - fh)/2
                            add_image_exact(sl, ci['stream'], sx, sy, fw1, fh)
                            add_image_exact(sl, ni['stream'], sx+fw1+G, sy, fw2, fh)
                            i += 2
                        else:
                            r = ci['w']/ci['h']
                            if uh*r <= uw:
                                fh, fw = uh, uh*r
                            else:
                                fh, fw = uw/r, uw
                                
                            align = av.get() if not isp else '2'
                                
                            if align == '1': sx = CL
                            elif align == '3': sx = sw - Inches(0.2) - fw
                            else: sx = CL + (uw - fw)/2
                                
                            sy = CT + (uh - fh)/2
                            add_image_exact(sl, ci['stream'], sx, sy, fw, fh)
                            i += 1
                        vt += 1
                else:
                    lms = [im for im in imd if not im['is_portrait']]
                    pms = [im for im in imd if im['is_portrait']]
                    
                    chunks = []
                    for c in partition_images(lms, 6): chunks.append({'t':'l', 'imgs':c})
                    for c in partition_images(pms, 4): chunks.append({'t':'p', 'imgs':c})
                        
                    for ck in chunks:
                        sl = prs.slides.add_slide(sly)
                        
                        if not isp and mbg_var.get() and os.path.exists(mbg_var.get()):
                            try:
                                sl.shapes.add_picture(mbg_var.get(), 0, 0, sw, sh)
                            except:
                                pass
                                
                        if not isp and t2_var.get():
                            p = sl.shapes.add_textbox(Inches(0.2), Inches(0.15), sw-Inches(0.4), Inches(0.6)).text_frame.paragraphs[0]
                            p.text = t2_var.get().upper()
                            p.font.size = Pt(22)
                            p.font.bold = True
                            p.font.underline = True
                            p.font.color.rgb = RGBColor(0, 51, 102)
                            
                        move_slide(prs, len(prs.slides)-1, vt)
                        
                        ims = ck['imgs']
                        n = len(ims)
                        if ck['t'] == 'p':
                            lrs = [ims]
                        else:
                            if n == 6: lrs = [ims[0:3], ims[3:6]]
                            elif n == 5: lrs = [ims[0:3], ims[3:5]]
                            elif n == 4: lrs = [ims[0:2], ims[2:4]]
                            else: lrs = [ims]
                                
                        draw_adaptive_grid(sl, lrs, CL, CT, uw, uh, G)
                        vt += 1
                        
                # Trang bìa kết thúc
                if not isp and (end_tit.get() or e_bg.get()):
                    add_cover(is_end=True)
                    
                prs.save(sp)
                messagebox.showinfo("OK", "Đã xuất PowerPoint thành công!")
                d.destroy()
            except PermissionError:
                messagebox.showerror("Lỗi", "File đang được mở, hãy tắt đi trước khi ghi đè!")
            except Exception as e:
                messagebox.showerror("Lỗi", f"PPT Error: {e}")

        def exp_pdf():
            pgs, fn = get_pgs()
            if not pgs: return
            sp = filedialog.asksaveasfilename(defaultextension=".pdf", filetypes=[("PDF", "*.pdf")], initialfile=f"{fn}.pdf")
            if not sp: return
            try:
                def mbg(img):
                    bg = Image.new("RGB", (3840, 2160), "white")
                    ir = img.width / img.height
                    if ir > 16/9:
                        nw, nh = 3840, int(3840 / ir)
                    else:
                        nw, nh = int(2160 * ir), 2160
                    bg.paste(img.resize((nw, nh), RESAMPLE_LANCZOS), ((3840-nw)//2, (2160-nh)//2))
                    return bg
                    
                imgs = [mbg(itm['rendered_image'].convert("RGB")) for itm in pgs]
                imgs[0].save(sp, save_all=True, append_images=imgs[1:])
                messagebox.showinfo("OK", "Xuất PDF thành công!")
                d.destroy()
            except PermissionError:
                messagebox.showerror("Lỗi", "File PDF đang mở, hãy tắt đi!")
            except Exception as e:
                messagebox.showerror("Lỗi", f"PDF Error: {e}")

        def exp_fld():
            pgs, fn = get_pgs()
            if not pgs: return
            fld = filedialog.askdirectory()
            if not fld: return
            try:
                for i, itm in enumerate(pgs):
                    itm['rendered_image'].convert("RGB").save(os.path.join(fld, f"{fn}_Page_{i+1}.jpg"), quality=100, subsampling=0)
                messagebox.showinfo("OK", "Lưu ảnh gốc thành công!")
                d.destroy()
            except PermissionError:
                messagebox.showerror("Lỗi", "Ảnh đang được mở, hãy tắt đi!")
            except Exception as e:
                messagebox.showerror("Lỗi", f"Folder Error: {e}")

        fb = tk.Frame(d)
        fb.pack(fill=tk.X, pady=20, padx=20)
        fb.columnconfigure(0, weight=1)
        fb.columnconfigure(1, weight=1)
        fb.columnconfigure(2, weight=1)
        
        tk.Button(fb, text="📊 Xuất PPTX", bg="#FF6600", fg="white", font=("Arial", 11, "bold"), command=exp_ppt).grid(row=0, column=0, padx=5, sticky="ew")
        tk.Button(fb, text="📄 Xuất PDF", bg="#DC3545", fg="white", font=("Arial", 11, "bold"), command=exp_pdf).grid(row=0, column=1, padx=5, sticky="ew")
        tk.Button(fb, text="📁 Lưu Ảnh Gốc", bg="#E0A800", fg="black", font=("Arial", 11, "bold"), command=exp_fld).grid(row=0, column=2, padx=5, sticky="ew")

    # --- HÀM THÊM ẢNH CHÈN ĐÈ ---
    def action_insert_image(self):
        if not self.original_image:
            messagebox.showwarning("Cảnh báo", "Vui lòng mở ảnh gốc trước!")
            return
        pth = filedialog.askopenfilename(filetypes=[("Image Files", "*.jpg;*.png;*.jpeg;*.webp")])
        if pth:
            try:
                img = Image.open(pth).convert("RGBA")
                if img.width > 500 or img.height > 500:
                    img.thumbnail((500, 500), RESAMPLE_LANCZOS)
                z = self.zoom_factor
                cx = self.canvas.canvasx(self.canvas.winfo_width()/2)/z
                cy = self.canvas.canvasy(self.canvas.winfo_height()/2)/z
                
                x1 = cx - img.width/2
                y1 = cy - img.height/2
                x2 = cx + img.width/2
                y2 = cy + img.height/2
                
                self.add_to_history("insert_image", (x1, y1, x2, y2), "", False)
                self.undo_stack[-1]['img_data'] = img
                self.current_tool.set("select")
                self.redraw_all()
            except Exception as e:
                messagebox.showerror("Lỗi", f"Không thể chèn ảnh: {e}")

    # --- CHỌN MÀU & HIỆU ỨNG CHỮ ---
    def choose_line_color(self):
        c = colorchooser.askcolor(title="Màu")[1]
        if c:
            self.line_color = c
            self.btn_color.config(bg=self.line_color)
            
    def choose_text_color(self):
        c = colorchooser.askcolor(title="Màu")[1]
        if c:
            self.text_color = c
            self.btn_text_color.config(bg=self.text_color)
            
    def choose_text_bg(self):
        c = colorchooser.askcolor(title="Nền")[1]
        if c:
            self.text_bg = c
            self.btn_text_bg.config(bg=self.text_bg)

    def apply_text_case(self, text):
        c = self.text_case.get()
        if c == "lowercase": return text.lower()
        if c == "UPPERCASE": return text.upper()
        if c == "Capitalize Each Word": return text.title()
        if c == "tOGGLE cASE": return text.swapcase()
        return text

    def add_to_history(self, action_type, coords, text, manual_text=False):
        self.action_counter += 1
        tag = f"action_{self.action_counter}"
        self.undo_stack.append({
            'tag': tag,
            'type': action_type,
            'coords': coords,
            'text': text,
            'manual_text': manual_text,
            'color': self.line_color,
            'width': self.line_width.get(),
            'font': self.font_family.get(),
            'size': self.font_size.get(),
            'bold': self.is_bold.get(),
            'italic': self.is_italic.get(),
            'underline': self.is_underline.get(),
            't_color': self.text_color,
            't_bg': self.text_bg
        })
        self.redo_stack.clear()
        return tag

    def undo(self, event=None):
        if self.undo_stack:
            self.redo_stack.append(self.undo_stack.pop())
            self.redraw_all()
            
    def redo(self, event=None):
        if self.redo_stack:
            self.undo_stack.append(self.redo_stack.pop())
            self.redraw_all()

    def bind_shortcuts(self):
        self.root.bind("<Control-z>", self.undo)
        self.root.bind("<Control-y>", self.redo)
        self.root.bind("<Control-s>", lambda e: self.add_to_right_panel())
        self.root.bind("<Control-a>", self.batch_edit_dialog)
        self.root.bind("<Control-A>", self.batch_edit_dialog)

    def batch_edit_dialog(self, event=None):
        if not self.undo_stack:
            messagebox.showwarning("Trống", "Không có đối tượng!")
            return
        nw = self.line_width.get()
        ns = self.font_size.get()
        nf = self.font_family.get()
        for act in self.undo_stack:
            act.update({
                'width': nw,
                'size': ns,
                'color': self.line_color,
                't_color': self.text_color,
                't_bg': self.text_bg,
                'font': nf,
                'bold': self.is_bold.get(),
                'italic': self.is_italic.get(),
                'underline': self.is_underline.get()
            })
        self.redraw_all()

    # --- HỆ THỐNG DANH SÁCH ẢNH TÍCH HỢP XÓA ---
    def add_multiple_images(self):
        pths = filedialog.askopenfilenames(filetypes=[("Image Files", "*.jpg;*.png;*.jpeg")])
        if not pths: return
        for p in pths:
            try:
                self.raw_images.append({'image': Image.open(p).convert("RGB")})
            except Exception:
                pass
        self.refresh_left_panel()
        if not self.original_image and self.raw_images:
            self.load_from_left(len(self.raw_images)-1)

    def refresh_left_panel(self):
        for w in self.left_scroll.scrollable_frame.winfo_children(): w.destroy()
        self.left_thumbs = []
        for i, itm in enumerate(self.raw_images):
            im = itm['image'].copy()
            im.thumbnail((240, 240))
            tk_im = ImageTk.PhotoImage(im)
            self.left_thumbs.append(tk_im)
            
            f = tk.Frame(self.left_scroll.scrollable_frame, bg="#252526")
            f.pack(pady=5, padx=5)
            
            tf = tk.Frame(f, bg="#252526")
            tf.pack(fill=tk.X)
            tk.Label(tf, text=f"Gốc {i+1}", bg="#252526", fg="white", font=("Arial", 9)).pack(side=tk.LEFT)
            tk.Button(tf, text="X", fg="red", bg="#252526", bd=0, command=lambda idx=i: self.delete_raw_image(idx)).pack(side=tk.RIGHT)
            
            lbl = tk.Label(f, image=tk_im, bg="#444", bd=2, cursor="hand2")
            lbl.pack()
            lbl.bind("<Button-1>", lambda e, idx=i: self.load_from_left(idx))
            
        tk.Button(self.left_scroll.scrollable_frame, text="➕ Thêm", font=("Arial", 12, "bold"), bg="#444", fg="white", relief="flat", command=self.add_multiple_images).pack(pady=10, padx=5, fill=tk.X)
        
        self.left_scroll.update_idletasks()
        self.left_scroll.canvas.configure(scrollregion=self.left_scroll.canvas.bbox("all"))

    def delete_raw_image(self, idx):
        if messagebox.askyesno("Xóa", "Xóa ảnh gốc này khỏi dự án?"):
            del self.raw_images[idx]
            if not self.raw_images:
                self.canvas.delete("all")
                self.original_image = None
                self.undo_stack.clear()
            self.refresh_left_panel()

    def refresh_right_panel(self):
        for w in self.right_scroll.scrollable_frame.winfo_children(): w.destroy()
        self.right_thumbs = []
        for i, itm in enumerate(self.exported_pages):
            im = itm['rendered_image'].copy()
            im.thumbnail((240, 240))
            tk_im = ImageTk.PhotoImage(im)
            self.right_thumbs.append(tk_im)
            
            f = tk.Frame(self.right_scroll.scrollable_frame, bg="#252526")
            f.pack(pady=5, padx=5)
            
            bc = "#28A745" if self.current_right_idx == i else "#444"
            
            tf = tk.Frame(f, bg="#252526")
            tf.pack(fill=tk.X)
            tk.Label(tf, text=f"Trang {i+1}", bg="#252526", fg="white", font=("Arial", 9)).pack(side=tk.LEFT)
            tk.Button(tf, text="X", fg="#E0A800", bg="#252526", bd=0, command=lambda idx=i: self.delete_export_image(idx)).pack(side=tk.RIGHT)
            
            lbl = tk.Label(f, image=tk_im, bg=bc, bd=3, cursor="hand2")
            lbl.pack()
            lbl.bind("<Double-Button-1>", lambda e, idx=i: self.load_from_right(idx))
            
        self.right_scroll.update_idletasks()
        self.right_scroll.canvas.configure(scrollregion=self.right_scroll.canvas.bbox("all"))

    def delete_export_image(self, idx):
        if messagebox.askyesno("Xóa", "Xóa ảnh này khỏi danh sách xuất?"):
            del self.exported_pages[idx]
            if self.current_right_idx == idx:
                self.current_right_idx = None
            self.refresh_right_panel()

    def load_from_left(self, idx):
        if idx >= len(self.raw_images): return
        self.canvas.delete("all")
        self.original_image = self.raw_images[idx]['image'].copy()
        self.undo_stack.clear()
        self.redo_stack.clear()
        self.current_right_idx = None
        self.zoom_factor = 1.0
        self.apply_zoom()
        self.refresh_right_panel() 

    def load_from_right(self, idx):
        if idx >= len(self.exported_pages): return
        self.canvas.delete("all")
        self.original_image = self.exported_pages[idx]['raw_image'].copy()
        self.undo_stack = copy.deepcopy(self.exported_pages[idx]['undo_stack'])
        self.redo_stack.clear()
        self.current_right_idx = idx
        self.zoom_factor = 1.0
        self.apply_zoom()
        self.refresh_right_panel() 

    def add_to_right_panel(self):
        if not self.original_image: return
        try:
            exp_img = self.get_export_image()
            if not exp_img: return
            
            safe_stack = []
            for act in self.undo_stack:
                na = {k: v for k, v in act.items() if k != 'tk_cache'}
                if na.get('img_data'):
                    na['img_data'] = na['img_data'].copy()
                safe_stack.append(copy.deepcopy(na))
                
            item = {
                'raw_image': self.original_image.copy(),
                'undo_stack': safe_stack,
                'rendered_image': exp_img
            }
            
            self.exported_pages.append(item)
            self.current_right_idx = len(self.exported_pages) - 1
            self.refresh_right_panel()
            
            self.canvas.config(bg="#28A745")
            self.root.after(100, lambda: self.canvas.config(bg="#1E1E1E"))
        except Exception as e:
            traceback.print_exc()
            messagebox.showerror("Lỗi Lưu", f"Lỗi: {e}")

    # --- HÀM VẼ LẠI TOÀN BỘ THEO Z-ORDER ---
    def redraw_all(self):
        self.canvas.delete("shape")
        self.canvas.delete("text_element")
        for act in self.undo_stack:
            self.render_action_on_canvas(act)
        try:
            for tg in ["layer_blur", "layer_insert_image", "layer_shape", "layer_dim", "layer_text"]:
                self.canvas.tag_raise(tg)
        except Exception:
            pass

    def on_tool_change(self, *args):
        t = self.current_tool.get()
        if t not in ["draw", "line", "rect", "oval", "blur"]:
            self.draw_btn.config(text="🎨 Các Hình Vẽ ▼")
            
        cur = {"select": "arrow", "draw": "pencil", "pan": "fleur"}.get(t, "crosshair")
        self.canvas.config(cursor=cur)
        self.redraw_all()

    def on_mouse_motion(self, event):
        if self.current_tool.get() != "select": return
        x, y = self.canvas.canvasx(event.x), self.canvas.canvasy(event.y)
        z = self.zoom_factor
        
        for act in reversed(self.undo_stack):
            t = act['type']
            if t in ["dim", "arrow_note", "line"]:
                x1, y1, x2, y2 = [v * z for v in act['coords']]
                if math.hypot(x-x1, y-y1)<15 or math.hypot(x-x2, y-y2)<15:
                    self.canvas.config(cursor="cross")
                    return
            elif t == "angle":
                if any(math.hypot(x-px*z, y-py*z)<15 for px, py in act['coords']):
                    self.canvas.config(cursor="cross")
                    return
            elif t == "draw":
                if any(math.hypot(x-px*z, y-py*z)<15 for px, py in act['coords'][::5]):
                    self.canvas.config(cursor="fleur")
                    return
            elif t in ["rect", "oval", "blur", "insert_image"]:
                if len(act['coords']) == 4:
                    x1, y1, x2, y2 = [v * z for v in act['coords']]
                    px1, py1, px2, py2 = min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)
                    
                    if t == "insert_image":
                        for hn, (hx, hy) in {
                            "scale_TL": (px1, py1), "scale_T": ((px1+px2)/2, py1), "scale_TR": (px2, py1),
                            "scale_L": (px1, (py1+py2)/2), "scale_R": (px2, (py1+py2)/2),
                            "scale_BL": (px1, py2), "scale_B": ((px1+px2)/2, py2), "scale_BR": (px2, py2)
                        }.items():
                            if abs(x-hx) <= 8 and abs(y-hy) <= 8:
                                cur_map = {
                                    "scale_TL":"size_nw_se", "scale_TR":"size_ne_sw", "scale_BL":"size_ne_sw", "scale_BR":"size_nw_se",
                                    "scale_L":"size_we", "scale_R":"size_we", "scale_T":"size_ns", "scale_B":"size_ns"
                                }
                                self.canvas.config(cursor=cur_map[hn])
                                return
                    else:
                        if math.hypot(x-px1, y-py1)<15 or math.hypot(x-px2, y-py2)<15:
                            self.canvas.config(cursor="size_nw_se")
                            return
                        if math.hypot(x-px2, y-py1)<15 or math.hypot(x-px1, y-py2)<15:
                            self.canvas.config(cursor="size_ne_sw")
                            return
                            
                    if px1 <= x <= px2 and py1 <= y <= py2:
                        self.canvas.config(cursor="fleur")
                        return
                        
        item = self.canvas.find_withtag("current")
        if item:
            tags = self.canvas.gettags(item[0])
            if "text_element" in tags:
                self.canvas.config(cursor="xterm")
                return
            elif "shape" in tags:
                self.canvas.config(cursor="fleur")
                return
        self.canvas.config(cursor="arrow")

    def render_action_on_canvas(self, act):
        z = self.zoom_factor
        t = act['type']
        tag = act['tag']
        color = act['color']
        lw = max(1, int(act['width'] * z))
        
        lt = "layer_shape"
        if t == "blur": lt = "layer_blur"
        elif t == "insert_image": lt = "layer_insert_image"
        elif t in ["dim", "arrow_note", "angle"]: lt = "layer_dim"

        if t == "dim":
            x1, y1, x2, y2 = [v * z for v in act['coords']]
            r = max(4, int(lw * 1.5)) 
            self.canvas.create_line(x1, y1, x2, y2, fill=color, width=lw, tags=("vector", tag, "shape", lt))
            self.canvas.create_oval(x1-r, y1-r, x1+r, y1+r, fill=color, outline=color, tags=("vector", tag, "shape", lt))
            self.canvas.create_oval(x2-r, y2-r, x2+r, y2+r, fill=color, outline=color, tags=("vector", tag, "shape", lt))
            self.draw_tk_text((x1+x2)/2, (y1+y2)/2, act['text'], tag, act)
            
        elif t == "line":
            self.canvas.create_line(*[v * z for v in act['coords']], fill=color, width=lw, tags=("vector", tag, "shape", lt))
            
        elif t == "rect":
            self.canvas.create_rectangle(*[v * z for v in act['coords']], outline=color, width=lw, tags=("vector", tag, "shape", lt))
            
        elif t == "oval":
            self.canvas.create_oval(*[v * z for v in act['coords']], outline=color, width=lw, tags=("vector", tag, "shape", lt))
            
        elif t == "blur":
            if self.original_image:
                rx1, ry1, rx2, ry2 = act['coords']
                l = int(max(0, min(rx1, rx2)))
                t_y = int(max(0, min(ry1, ry2)))
                r = int(min(self.original_image.width, max(rx1, rx2)))
                b = int(min(self.original_image.height, max(ry1, ry2)))
                
                wc = r - l
                hc = b - t_y
                
                if wc > 0 and hc > 0:
                    ps = 40
                    sw = max(1, wc // ps)
                    sh = max(1, hc // ps)
                    blur = self.original_image.crop((l, t_y, r, b)).resize((sw, sh), RESAMPLE_BOX).resize((wc, hc), RESAMPLE_NEAREST)
                    
                    nw = int(wc * z)
                    nh = int(hc * z)
                    if nw > 0 and nh > 0:
                        t_im = ImageTk.PhotoImage(blur.resize((nw, nh), RESAMPLE_NEAREST))
                        act['tk_cache'] = t_im 
                        self.canvas.create_image(min(act['coords'][0]*z, act['coords'][2]*z), min(act['coords'][1]*z, act['coords'][3]*z), anchor=tk.NW, image=t_im, tags=("vector", tag, "shape", lt))
                        if self.current_tool.get() == "select":
                            self.canvas.create_rectangle(min(act['coords'][0]*z, act['coords'][2]*z), min(act['coords'][1]*z, act['coords'][3]*z), max(act['coords'][0]*z, act['coords'][2]*z), max(act['coords'][1]*z, act['coords'][3]*z), outline="#333", dash=(4,4), width=1, tags=("vector", tag, "shape", lt))
                            
        elif t == "insert_image":
            x1, y1, x2, y2 = [v * z for v in act['coords']]
            nw = int(abs(x2-x1))
            nh = int(abs(y2-y1))
            if nw > 0 and nh > 0:
                t_im = ImageTk.PhotoImage(act['img_data'].resize((nw, nh), RESAMPLE_LANCZOS))
                act['tk_cache'] = t_im
                self.canvas.create_image(min(x1, x2), min(y1, y2), anchor=tk.NW, image=t_im, tags=("vector", tag, "shape", lt))
                if self.current_tool.get() == "select":
                    px1, py1, px2, py2 = min(x1, x2), min(y1, y2), max(x1, x2), max(y1, y2)
                    self.canvas.create_rectangle(px1, py1, px2, py2, outline="#007ACC", width=2, tags=("vector", tag, "shape", lt))
                    for hx, hy in [(px1, py1), ((px1+px2)/2, py1), (px2, py1), (px1, (py1+py2)/2), (px2, (py1+py2)/2), (px1, py2), ((px1+px2)/2, py2), (px2, py2)]:
                        self.canvas.create_rectangle(hx-5, hy-5, hx+5, hy+5, fill="white", outline="#007ACC", width=2, tags=("vector", tag, "shape", lt))
                        
        elif t == "arrow_note":
            x1, y1, x2, y2 = [v * z for v in act['coords']]
            self.canvas.create_line(x1, y1, x2, y2, fill=color, width=lw, arrow=tk.LAST, arrowshape=(max(15, int(lw*4)), max(18, int(lw*5)), max(6, int(lw*2.5))), tags=("vector", tag, "shape", lt))
            self.draw_tk_text(x1, y1, act['text'], tag, act)
            
        elif t == "tag":
            self.draw_tk_text(*[v * z for v in act['coords']], act['text'], tag, act)
            
        elif t == "angle":
            pts = [(px*z, py*z) for px, py in act['coords']]
            p1, p2, p3 = pts
            d1 = math.degrees(math.atan2(-(p2[1]-p1[1]), p2[0]-p1[0])) % 360
            d2 = math.degrees(math.atan2(-(p3[1]-p1[1]), p3[0]-p1[0])) % 360
            df = (d2 - d1) % 360
            if df > 180:
                d1, d2, df = d2, d1, 360 - df
            al1, al2, aw = max(15, int(lw*4)), max(18, int(lw*5)), max(6, int(lw*2.5))
            self.canvas.create_line(p1[0], p1[1], p2[0], p2[1], fill=color, width=lw, arrow=tk.LAST, arrowshape=(al1, al2, aw), tags=("vector", tag, "shape", lt))
            self.canvas.create_line(p1[0], p1[1], p3[0], p3[1], fill=color, width=lw, arrow=tk.LAST, arrowshape=(al1, al2, aw), tags=("vector", tag, "shape", lt))
            r = min(35*z, math.hypot(p2[0]-p1[0], p2[1]-p1[1])*0.8, math.hypot(p3[0]-p1[0], p3[1]-p1[1])*0.8)
            if df > 0 and r > 5:
                self.canvas.create_arc(p1[0]-r, p1[1]-r, p1[0]+r, p1[1]+r, start=d1, extent=df, style=tk.ARC, outline=color, width=max(1, int(lw*0.8)), tags=("vector", tag, "shape", lt))
            md = d1 + df/2
            tx = p1[0] + (r+20*z) * math.cos(math.radians(md))
            ty = p1[1] - (r+20*z) * math.sin(math.radians(md))
            self.draw_tk_text(tx, ty, act['text'], tag, act)
            
        elif t == "draw":
            fpts = [cv * z for pt in act['coords'] for cv in pt]
            if len(fpts) >= 4:
                self.canvas.create_line(*fpts, fill=color, width=lw, capstyle=tk.ROUND, joinstyle=tk.ROUND, smooth=True, tags=("vector", tag, "shape", lt))
            elif len(fpts) == 2:
                self.canvas.create_oval(fpts[0]-lw/2, fpts[1]-lw/2, fpts[0]+lw/2, fpts[1]+lw/2, fill=color, outline=color, tags=("vector", tag, "shape", lt))

    def draw_tk_text(self, x, y, text, action_tag, act_dict=None):
        if not act_dict.get('text', ''): return
        fs = max(4, int(act_dict['size'] * self.zoom_factor))
        font = tkfont.Font(family=act_dict['font'], size=fs, weight="bold" if act_dict['bold'] else "normal", slant="italic" if act_dict['italic'] else "roman", underline=act_dict['underline'])
        tmp = self.canvas.create_text(x, y, text=act_dict['text'], font=font)
        bbox = self.canvas.bbox(tmp)
        self.canvas.delete(tmp)
        p = max(2, int(6 * self.zoom_factor))
        self.canvas.create_rectangle(bbox[0]-p, bbox[1]-p, bbox[2]+p, bbox[3]+p, fill=act_dict['t_bg'], outline=act_dict['t_bg'], tags=("vector", action_tag, "text_element", "layer_text"))
        self.canvas.create_text(x, y, text=act_dict['text'], font=font, fill=act_dict['t_color'], tags=("vector", action_tag, "text_element", "layer_text"))

    def on_double_click(self, event):
        if self.current_tool.get() != "select": return
        item = self.canvas.find_withtag("current")
        if item:
            st = next((t for t in self.canvas.gettags(item[0]) if t.startswith("action_")), None)
            if st:
                for act in self.undo_stack:
                    if act['tag'] == st and act['type'] not in ["draw", "blur", "insert_image"]:
                        self.open_edit_properties_dialog(act)
                        break

    def open_edit_properties_dialog(self, act):
        self.show_virtual_keyboard() 
        d = tk.Toplevel(self.root)
        d.title("Sửa")
        d.geometry("+%d+%d" % (self.root.winfo_rootx()+350, self.root.winfo_rooty()+200))
        d.transient(self.root)
        d.grab_set()
        
        t_var = tk.StringVar(value=act.get('text', ''))
        tk.Entry(d, textvariable=t_var, width=40, font=("Arial", 11)).pack(pady=15)
        
        fp = tk.Frame(d)
        fp.pack(pady=5)
        s_var = tk.IntVar(value=act['size'])
        tk.Spinbox(fp, from_=8, to=200, textvariable=s_var, width=10).grid(row=0, column=1)
        w_var = tk.IntVar(value=act['width'])
        tk.Spinbox(fp, from_=1, to=50, textvariable=w_var, width=10).grid(row=1, column=1)
        
        def save():
            act.update({'text': t_var.get(), 'size': s_var.get(), 'width': w_var.get(), 'manual_text': True})
            self.redraw_all()
            d.destroy()
            
        tk.Button(d, text="✅ Áp dụng", command=save, bg="#28A745", fg="white").pack(pady=15)

    def on_left_press(self, event):
        if not self.original_image: return
        x, y = self.canvas.canvasx(event.x), self.canvas.canvasy(event.y)
        z = self.zoom_factor
        t = self.current_tool.get()
        
        if t == "pan":
            self.canvas.scan_mark(event.x, event.y)
            return
            
        if t == "select":
            self.drag_mode = None
            for act in reversed(self.undo_stack):
                at = act['type']
                if at in ["dim", "arrow_note", "line"]:
                    x1, y1, x2, y2 = [v * z for v in act['coords']]
                    if math.hypot(x-x1, y-y1) < 15:
                        self.selected_action, self.drag_mode = act['tag'], "move_p1"
                        return
                    if math.hypot(x-x2, y-y2) < 15:
                        self.selected_action, self.drag_mode = act['tag'], "move_p2"
                        return
                    if at == "line" and min(x1, x2)-15 <= x <= max(x1, x2)+15 and min(y1, y2)-15 <= y <= max(y1, y2)+15:
                        self.selected_action, self.drag_mode = act['tag'], "move_all"
                        self.last_drag_x, self.last_drag_y = x, y
                        self.total_dx, self.total_dy = 0, 0
                        return
                elif at == "angle":
                    pts = [(px*z, py*z) for px, py in act['coords']]
                    if math.hypot(x-pts[0][0], y-pts[0][1]) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p1"; return
                    if math.hypot(x-pts[1][0], y-pts[1][1]) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p2"; return
                    if math.hypot(x-pts[2][0], y-pts[2][1]) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p3"; return
                elif at in ["rect", "oval", "blur", "insert_image"]:
                    px1, py1, px2, py2 = min(act['coords'][0]*z, act['coords'][2]*z), min(act['coords'][1]*z, act['coords'][3]*z), max(act['coords'][0]*z, act['coords'][2]*z), max(act['coords'][1]*z, act['coords'][3]*z)
                    if at == "insert_image":
                        handles = {
                            "scale_TL": (px1, py1), "scale_T": ((px1+px2)/2, py1), "scale_TR": (px2, py1),
                            "scale_L": (px1, (py1+py2)/2), "scale_R": (px2, (py1+py2)/2),
                            "scale_BL": (px1, py2), "scale_B": ((px1+px2)/2, py2), "scale_BR": (px2, py2)
                        }
                        for hn, (hx, hy) in handles.items():
                            if abs(x-hx) <= 10 and abs(y-hy) <= 10:
                                self.selected_action, self.drag_mode = act['tag'], hn
                                self.drag_start_rect = [act['coords'][0]*z, act['coords'][1]*z, act['coords'][2]*z, act['coords'][3]*z]
                                self.drag_start_x, self.drag_start_y = x, y
                                return
                    else:
                        if math.hypot(x-px1, y-py1) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p1"; return
                        if math.hypot(x-px2, y-py2) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p2"; return
                        if math.hypot(x-px2, y-py1) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p3"; return
                        if math.hypot(x-px1, y-py2) < 15: self.selected_action, self.drag_mode = act['tag'], "move_p4"; return
                    if px1 <= x <= px2 and py1 <= y <= py2:
                        self.selected_action, self.drag_mode = act['tag'], "move_all"
                        self.last_drag_x, self.last_drag_y = x, y
                        self.total_dx, self.total_dy = 0, 0
                        return
                        
            item = self.canvas.find_withtag("current")
            if item:
                st = next((tg for tg in self.canvas.gettags(item[0]) if tg.startswith("action_")), None)
                if st:
                    self.selected_action, self.drag_mode = st, "move_all"
                    self.last_drag_x, self.last_drag_y = x, y
                    self.total_dx, self.total_dy = 0, 0
            else:
                self.selected_action = None
                
        elif t == "dim":
            self.start_x, self.start_y = x, y
            lw = max(1, int(self.line_width.get()*z))
            r = max(4, int(lw*1.5)) 
            self.temp_shapes = [
                self.canvas.create_line(x, y, x, y, fill=self.line_color, width=lw),
                self.canvas.create_oval(x-r, y-r, x+r, y+r, fill=self.line_color, outline=self.line_color),
                self.canvas.create_oval(x-r, y-r, x+r, y+r, fill=self.line_color, outline=self.line_color)
            ]
        elif t in ["line", "rect", "oval", "blur"]:
            self.start_x, self.start_y = x, y
            lw = max(1, int(self.line_width.get()*z))
            if t == "line": self.temp_shapes = [self.canvas.create_line(x, y, x, y, fill=self.line_color, width=lw)]
            elif t == "rect": self.temp_shapes = [self.canvas.create_rectangle(x, y, x, y, outline=self.line_color, width=lw)]
            elif t == "oval": self.temp_shapes = [self.canvas.create_oval(x, y, x, y, outline=self.line_color, width=lw)]
            elif t == "blur": self.temp_shapes = [self.canvas.create_rectangle(x, y, x, y, outline="#007ACC", dash=(4,4), width=2)]
        elif t == "arrow_note":
            self.start_x, self.start_y = x, y
            lw = max(1, int(self.line_width.get()*z))
            self.temp_shapes = [self.canvas.create_line(x, y, x, y, fill=self.line_color, width=lw, arrow=tk.LAST, arrowshape=(max(15, int(lw*4)), max(18, int(lw*5)), max(6, int(lw*2.5))))]
        elif t == "tag":
            self.show_virtual_keyboard()
            r_v = simpledialog.askstring("Note", "Nhập ghi chú:")
            if r_v:
                self.add_to_history("tag", (x/z, y/z), self.apply_text_case(r_v), True)
                self.redraw_all()
        elif t == "angle":
            self.angle_points.append((x, y))
            self.temp_shapes.append(self.canvas.create_oval(x-3, y-3, x+3, y+3, fill=self.line_color))
            if len(self.angle_points) == 3:
                p1, p2, p3 = self.angle_points
                d = abs(math.degrees(math.atan2(p3[1]-p1[1], p3[0]-p1[0])) - math.degrees(math.atan2(p2[1]-p1[1], p2[0]-p1[0])))
                if d > 180: d = 360 - d
                final_text = f"{d:.1f}°"
                self.add_to_history("angle", ((p1[0]/z, p1[1]/z), (p2[0]/z, p2[1]/z), (p3[0]/z, p3[1]/z)), final_text, False)
                for s in self.temp_shapes: self.canvas.delete(s)
                self.temp_shapes, self.angle_points = [], []
                self.redraw_all()
        elif t == "draw":
            self.draw_points = [(x, y)]
            lw = max(1, int(self.line_width.get()*z))
            self.temp_shapes.append(self.canvas.create_oval(x-lw/2, y-lw/2, x+lw/2, y+lw/2, fill=self.line_color, outline=self.line_color))

    def on_left_drag(self, event):
        t = self.current_tool.get()
        cx = self.canvas.canvasx(event.x)
        cy = self.canvas.canvasy(event.y)
        
        if t == "pan":
            self.canvas.scan_dragto(event.x, event.y, gain=1)
            return
            
        if t == "select" and self.selected_action:
            if self.drag_mode == "move_all":
                dx = cx - self.last_drag_x
                dy = cy - self.last_drag_y
                self.canvas.move(self.selected_action, dx, dy)
                self.last_drag_x, self.last_drag_y = cx, cy
                self.total_dx += dx
                self.total_dy += dy
                
            elif self.drag_mode and self.drag_mode.startswith("scale_"):
                z = self.zoom_factor
                for act in self.undo_stack:
                    if act['tag'] == self.selected_action:
                        dx = cx - self.drag_start_x
                        dy = cy - self.drag_start_y
                        nx1, ny1, nx2, ny2 = self.drag_start_rect
                        
                        if "L" in self.drag_mode: nx1 += dx
                        if "R" in self.drag_mode: nx2 += dx
                        if "T" in self.drag_mode: ny1 += dy
                        if "B" in self.drag_mode: ny2 += dy
                        
                        if nx2 - nx1 < 20*z:
                            if "L" in self.drag_mode: nx1 = nx2 - 20*z
                            if "R" in self.drag_mode: nx2 = nx1 + 20*z
                        if ny2 - ny1 < 20*z:
                            if "T" in self.drag_mode: ny1 = ny2 - 20*z
                            if "B" in self.drag_mode: ny2 = ny1 + 20*z
                            
                        if event.state & 0x0001: 
                            ir = act['img_data'].width / max(1.0, act['img_data'].height)
                            nw = max(1.0, nx2 - nx1)
                            nh = max(1.0, ny2 - ny1)
                            if nw / nh > ir: nh = nw / ir
                            else: nw = nh * ir
                            
                            if self.drag_mode in ["scale_L", "scale_R"]:
                                my = (ny1 + ny2)/2
                                ny1, ny2 = my - nh/2, my + nh/2
                            elif self.drag_mode in ["scale_T", "scale_B"]:
                                mx = (nx1 + nx2)/2
                                nx1, nx2 = mx - nw/2, mx + nw/2
                            else:
                                if "L" in self.drag_mode: nx1 = nx2 - nw
                                if "R" in self.drag_mode: nx2 = nx1 + nw
                                if "T" in self.drag_mode: ny1 = ny2 - nh
                                if "B" in self.drag_mode: ny2 = ny1 + nh
                                
                        act['coords'] = (nx1/z, ny1/z, nx2/z, ny2/z)
                        self.redraw_all()
                        break
                        
            elif "move_p" in self.drag_mode:
                z = self.zoom_factor
                for act in self.undo_stack:
                    if act['tag'] == self.selected_action:
                        rx, ry = cx/z, cy/z
                        at = act['type']
                        if at in ["dim", "arrow_note", "line"]:
                            act['coords'] = (rx, ry, act['coords'][2], act['coords'][3]) if self.drag_mode == "move_p1" else (act['coords'][0], act['coords'][1], rx, ry)
                        elif at == "angle":
                            p1, p2, p3 = act['coords']
                            if self.drag_mode == "move_p1": p1 = (rx, ry)
                            elif self.drag_mode == "move_p2": p2 = (rx, ry)
                            elif self.drag_mode == "move_p3": p3 = (rx, ry)
                            act['coords'] = (p1, p2, p3)
                            if not act.get('manual_text'):
                                d = abs(math.degrees(math.atan2(p3[1]-p1[1], p3[0]-p1[0])) - math.degrees(math.atan2(p2[1]-p1[1], p2[0]-p1[0])))
                                if d > 180: d = 360 - d
                                act['text'] = f"{d:.1f}°"
                        elif at in ["rect", "oval", "blur"]:
                            x1, y1, x2, y2 = act['coords']
                            if x1 > x2: x1, x2 = x2, x1
                            if y1 > y2: y1, y2 = y2, y1
                            if self.drag_mode == "move_p1": act['coords'] = (rx, ry, x2, y2)
                            elif self.drag_mode == "move_p2": act['coords'] = (x1, y1, rx, ry)
                            elif self.drag_mode == "move_p3": act['coords'] = (x1, ry, rx, y2)
                            elif self.drag_mode == "move_p4": act['coords'] = (rx, y1, x2, ry)
                        self.redraw_all()
                        break
                        
        elif t in ["dim", "arrow_note", "line", "rect", "oval", "blur"] and self.temp_shapes:
            self.canvas.coords(self.temp_shapes[0], self.start_x, self.start_y, cx, cy)
            if t == "dim":
                lw = max(1, int(self.line_width.get() * self.zoom_factor))
                r = max(4, int(lw * 1.5)) 
                self.canvas.coords(self.temp_shapes[1], self.start_x-r, self.start_y-r, self.start_x+r, self.start_y+r)
                self.canvas.coords(self.temp_shapes[2], cx-r, cy-r, cx+r, cy+r)
                
        elif t == "draw" and self.draw_points:
            self.draw_points.append((cx, cy))
            lw = max(1, int(self.line_width.get() * self.zoom_factor))
            self.temp_shapes.append(self.canvas.create_line(self.draw_points[-2][0], self.draw_points[-2][1], cx, cy, fill=self.line_color, width=lw, capstyle=tk.ROUND, smooth=True))

    def on_left_release(self, event):
        t = self.current_tool.get()
        z = self.zoom_factor
        
        if t == "pan": return
        if t == "draw":
            if self.draw_points:
                self.add_to_history("draw", [(px/z, py/z) for px, py in self.draw_points], "", False)
                for s in self.temp_shapes: self.canvas.delete(s)
                self.temp_shapes = []
                self.draw_points = []
                self.redraw_all()
            return
            
        if t == "select" and self.selected_action:
            if self.drag_mode == "move_all":
                rdx = self.total_dx / z
                rdy = self.total_dy / z
                for act in self.undo_stack:
                    if act['tag'] == self.selected_action:
                        at = act['type']
                        if at in ["dim", "arrow_note", "line", "rect", "oval", "blur", "insert_image"]:
                            act['coords'] = (act['coords'][0]+rdx, act['coords'][1]+rdy, act['coords'][2]+rdx, act['coords'][3]+rdy)
                        elif at in ["tag"]:
                            act['coords'] = (act['coords'][0]+rdx, act['coords'][1]+rdy)
                        elif at == "angle":
                            act['coords'] = ((act['coords'][0][0]+rdx, act['coords'][0][1]+rdy), (act['coords'][1][0]+rdx, act['coords'][1][1]+rdy), (act['coords'][2][0]+rdx, act['coords'][2][1]+rdy))
                        elif at == "draw":
                            act['coords'] = [(px+rdx, py+rdy) for px, py in act['coords']]
                        break
            self.selected_action = None
            self.drag_mode = None
            return
            
        if not self.temp_shapes or t not in ["dim", "arrow_note", "line", "rect", "oval", "blur"]: return
        
        cx = self.canvas.canvasx(event.x)
        cy = self.canvas.canvasy(event.y)
        
        if abs(cx - self.start_x) < 5 and abs(cy - self.start_y) < 5:
            for s in self.temp_shapes: self.canvas.delete(s)
            self.temp_shapes = []
            return
            
        for s in self.temp_shapes: self.canvas.delete(s)
        self.temp_shapes = []
        
        if t in ["line", "rect", "oval", "blur"]:
            self.add_to_history(t, (self.start_x/z, self.start_y/z, cx/z, cy/z), "", False)
            self.redraw_all()
            return
            
        self.show_virtual_keyboard() 
        r_v = simpledialog.askstring("Nhập liệu", "Nhập thông số/ghi chú:")
        if r_v:
            if t == "dim" and self.dim_unit.get().strip() and not r_v.endswith(self.dim_unit.get().strip()):
                r_v = f"{r_v} {self.dim_unit.get().strip()}".strip()
            self.add_to_history(t, (self.start_x/z, self.start_y/z, cx/z, cy/z), self.apply_text_case(r_v), True)
            self.redraw_all()

    def get_export_image(self):
        if not self.original_image: return None
        exp = self.original_image.copy()
        dw = ImageDraw.Draw(exp)
        
        def gz(at):
            if at == "blur": return 0
            if at == "insert_image": return 1
            if at in ["draw", "line", "rect", "oval"]: return 2
            return 3
            
        for act in sorted(self.undo_stack, key=lambda a: gz(a['type'])):
            c = act['color']
            lw = max(1, int(act['width']*2))
            t = act['type']
            
            fn = "times.ttf" if "Times" in act['font'] else ("cour.ttf" if "Courier" in act['font'] else "arial.ttf")
            if act['bold'] and act['italic']: fn = fn.replace(".ttf", "z.ttf" if "times" not in fn else "bi.ttf")
            elif act['bold']: fn = fn.replace(".ttf", "bd.ttf")
            elif act['italic']: fn = fn.replace(".ttf", "i.ttf")
            try: pf = ImageFont.truetype(fn, max(8, int(act['size']*2.66)))
            except Exception: pf = ImageFont.load_default()

            if t == "dim":
                x1, y1, x2, y2 = act['coords']
                r = max(6, int(lw*1.5)) 
                dw.line([(x1, y1), (x2, y2)], fill=c, width=lw)
                dw.ellipse([x1-r, y1-r, x1+r, y1+r], fill=c)
                dw.ellipse([x2-r, y2-r, x2+r, y2+r], fill=c)
                self.dr_ptx(dw, act, (x1+x2)/2, (y1+y2)/2, pf)
            elif t == "line":
                dw.line([(act['coords'][0], act['coords'][1]), (act['coords'][2], act['coords'][3])], fill=c, width=lw)
            elif t == "rect":
                dw.rectangle([min(act['coords'][0], act['coords'][2]), min(act['coords'][1], act['coords'][3]), max(act['coords'][0], act['coords'][2]), max(act['coords'][1], act['coords'][3])], outline=c, width=lw)
            elif t == "oval":
                dw.ellipse([min(act['coords'][0], act['coords'][2]), min(act['coords'][1], act['coords'][3]), max(act['coords'][0], act['coords'][2]), max(act['coords'][1], act['coords'][3])], outline=c, width=lw)
            elif t == "blur":
                rx1, ry1, rx2, ry2 = act['coords']
                l = int(max(0, min(rx1, rx2)))
                t_y = int(max(0, min(ry1, ry2)))
                r = int(min(exp.width, max(rx1, rx2)))
                b = int(min(exp.height, max(ry1, ry2)))
                
                wc = r - l
                hc = b - t_y
                
                if wc > 0 and hc > 0:
                    ps = 40
                    sw = max(1, wc // ps)
                    sh = max(1, hc // ps)
                    blur = exp.crop((l, t_y, r, b)).resize((sw, sh), RESAMPLE_BOX).resize((wc, hc), RESAMPLE_NEAREST)
                    exp.paste(blur, (l, t_y, r, b))
                    
            elif t == "insert_image":
                x1, y1, x2, y2 = act['coords']
                nw = int(abs(x2-x1))
                nh = int(abs(y2-y1))
                if nw > 0 and nh > 0:
                    rs = act['img_data'].resize((nw, nh), RESAMPLE_LANCZOS)
                    if rs.mode == 'RGBA':
                        exp.paste(rs, (int(min(x1, x2)), int(min(y1, y2))), rs)
                    else:
                        exp.paste(rs, (int(min(x1, x2)), int(min(y1, y2))))
            elif t == "arrow_note":
                x1, y1, x2, y2 = act['coords']
                dw.line([(x1, y1), (x2, y2)], fill=c, width=lw)
                a = math.atan2(y1-y2, x1-x2)
                al = max(25, lw*5)
                aa = math.pi/6
                dw.polygon([(x2, y2), (x2+al*math.cos(a-aa), y2+al*math.sin(a-aa)), (x2+al*math.cos(a+aa), y2+al*math.sin(a+aa))], fill=c)
                self.dr_ptx(dw, act, x1, y1, pf)
            elif t == "tag":
                self.dr_ptx(dw, act, act['coords'][0], act['coords'][1], pf)
            elif t == "angle":
                p1, p2, p3 = act['coords']
                dw.line([p1, p2], fill=c, width=lw)
                dw.line([p1, p3], fill=c, width=lw)
                for ep in [p2, p3]:
                    a = math.atan2(p1[1]-ep[1], p1[0]-ep[0])
                    al = max(25, lw*5)
                    aa = math.pi/6
                    dw.polygon([ep, (ep[0]+al*math.cos(a-aa), ep[1]+al*math.sin(a-aa)), (ep[0]+al*math.cos(a+aa), ep[1]+al*math.sin(a+aa))], fill=c)
                d1 = math.degrees(math.atan2(-(p2[1]-p1[1]), p2[0]-p1[0])) % 360
                d2 = math.degrees(math.atan2(-(p3[1]-p1[1]), p3[0]-p1[0])) % 360
                df = (d2-d1) % 360
                if df > 180:
                    d1, d2, df = d2, d1, 360-df
                r = min(35, math.hypot(p2[0]-p1[0], p2[1]-p1[1])*0.8, math.hypot(p3[0]-p1[0], p3[1]-p1[1])*0.8)
                if df > 0 and r > 5:
                    dw.arc([p1[0]-r, p1[1]-r, p1[0]+r, p1[1]+r], 360-(d1+df), 360-d1, fill=c, width=max(1, int(lw*0.8)))
                md = d1+df/2
                self.dr_ptx(dw, act, p1[0]+(r+20)*math.cos(math.radians(md)), p1[1]-(r+20)*math.sin(math.radians(md)), pf)
            elif t == "draw":
                pts = [tuple(pt) for pt in act['coords']]
                if len(pts) > 1:
                    dw.line(pts, fill=c, width=lw, joint="curve")
                elif len(pts) == 1:
                    dw.ellipse([pts[0][0]-lw/2, pts[0][1]-lw/2, pts[0][0]+lw/2, pts[0][1]+lw/2], fill=c)
        return exp

    def dr_ptx(self, dw, act, cx, cy, f):
        if not act.get('text', ''): return
        bb = dw.textbbox((0, 0), act['text'], font=f)
        tw = bb[2]-bb[0]
        th = bb[3]-bb[1]
        p = max(4, int(th*0.25))
        rl = cx-tw/2-p
        rt = cy-th/2-p
        rr = cx+tw/2+p
        rb = cy+th/2+p
        dw.rectangle([rl, rt, rr, rb], fill=act['t_bg'])
        dw.text((cx-tw/2-bb[0], cy-th/2-bb[1]), act['text'], fill=act['t_color'], font=f)
        if act['underline']:
            dw.line([(rl+p, rb-p/2), (rr-p, rb-p/2)], fill=act['t_color'], width=max(1, int(act['size']*0.15)))

    def update_image_display(self):
        if not self.original_image: return
        nw = int(self.original_image.width * self.zoom_factor)
        nh = int(self.original_image.height * self.zoom_factor)
        self.tk_image = ImageTk.PhotoImage(self.original_image.resize((nw, nh), RESAMPLE_LANCZOS))
        self.canvas.create_image(0, 0, anchor=tk.NW, image=self.tk_image, tags="bg_image")
        self.canvas.tag_lower("bg_image")

    def on_right_press(self, event):
        self.canvas.scan_mark(event.x, event.y)
        
    def on_right_drag(self, event):
        self.canvas.scan_dragto(event.x, event.y, gain=1)
        
    def on_mousewheel(self, event):
        if not self.original_image: return
        if event.delta > 0:
            self.zoom_factor *= 1.1
        else:
            self.zoom_factor *= 0.9
        self.apply_zoom()

if __name__ == "__main__":
    root = tk.Tk()
    def handle_exception(exc_type, exc_value, exc_traceback):
        error_msg = "".join(traceback.format_exception(exc_type, exc_value, exc_traceback))
        print(error_msg)
        try:
            messagebox.showerror("Báo cáo lỗi (Crash Report)", f"Chi tiết lỗi:\n\n{error_msg}")
        except Exception:
            pass
    root.report_callback_exception = handle_exception
    sys.excepthook = handle_exception
    app = AdvancedNoteDimApp(root)
    root.mainloop()